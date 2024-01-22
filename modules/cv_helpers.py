import streamlit as st
import pdfplumber
import base64
import itertools
import base64
from PyPDF2 import PdfReader
from modules.openai_utils import *
from modules.app_helpers import *
from dotenv import load_dotenv
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
        text += "\n"
    return text

def extract_text_from_pdf_pypdf2(source_file):
    reader = PdfReader(source_file)
    text = ""
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text += page.extract_text() + "\n"  
    
    # Remove special characters from text
    text = ''.join(e for e in text if e.isalnum() or e.isspace())
    return text

def extract_text_from_pdf(source_file):
    with pdfplumber.open(source_file) as pdf:
        text = ""
        for page in pdf.pages:
            words = page.extract_words()

            # Sort words by their top and left coordinates
            sorted_words = sorted(words, key=lambda word: (-word['top'], word['x0']))

            # Group words based on their top coordinates
            grouped_words = itertools.groupby(sorted_words, key=lambda word: word['top'])

            # Combine words in each group, separated by spaces
            lines = []
            for _, word_group in grouped_words:
                line = ' '.join(word['text'] for word in word_group)
                lines.append(line)

            # Combine lines into a single string, separated by newlines
            page_text = '\n'.join(lines)
            text += page_text + "\n\n"
            
    text = ''.join(e for e in text if e.isalnum() or e.isspace())
    return text

def image_to_base64(image_path):
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode("utf-8")   
    
# Create a function that generates a prompt for OpenAI's GPT model based on the selected CV function
def generate_prompt_cv_function(cv_function, cv_custom_prompt=""):
    prompt = ""
    if cv_custom_prompt != "":
        prompt = cv_custom_prompt + "\n\n"
    else:
        if cv_function == "CV Reviewer":       
            # Update this prompt to include the CV Reviewer instructions.
            prompt = (
                "You are part of the recruitment team and your job is to review CVs of candidates\n"   
                "Respond with feedbacks and additional recommendations about contents of the CV below.\n"             
                "The owner of the CV is a Sofware Engineer with various experience in the IT industry.\n"
                "Recommend brevity to make the CV easier to read and more appealing to the readers. Choose 3 areas found in the CV and rewrite them to improve the content.\n"
                "Do not use Markdown or HTML tags in your response.\n"
                "State if the CV is well-written or not based on the following criteria:\n" 
                "<Start criteria>\n"       
                " 1. The CV should be clear and concise.\n"
                " 2. The CV should highlight the most important and relevant achievements.\n"
                " 3. The CV should be written in English.\n"
                " 4. The CV should have a Professional Background.\n"
                " 5. The CV should be written in a professional tone.\n"        
                " 6. The CV should contain at least 1 IT relevant experience.\n"
                " 7. The CV should contain at least 2 IT relevant skills.\n"
                " 8. The CV should contain at least 1 IT relevant certification.\n"   
                " 9. The CV should contain any prior experience, even if they're not IT related.\n"  
                " 10. The CV should contain educational achievements and awards.\n"
                "<End criteria>")
        elif cv_function == "CV Summarizer":
            prompt = (
                "You are part of the recruitment team and your job is to review CVs of candidates.\n"
                "Your job is to review the contents of the CV and generate a brief summary of the cv contents.\n"
                "Do not use Markdown or HTML tags in your response.\n"
                "Example of response:\n"
                "# Example:\n"
                "Summary of Professional Background [in less than 50 words]:\n"
                "The candidate is a Software Engineer with 10 years of experience in the IT industry.\n"
                "Summary of Technical Skills:\n"
                "1. Skill 1 \n"
                "2. Skill 2 \n"
                "3. Skill 3 \n"
                "Summary of Certifications:\n"
                "1. Certifications 1 \n"
                "2. Certifications 2 \n"
                "3. Certifications 3 \n"
                "Summary of Experiences:\n"
                "1. Experience 1 [less than 50 words] \n"
                "2. Experience 2 [less than 50 words] \n"
                "3. Experience 3 [less than 50 words] \n"
                "Other Notable mentions:\n"
                "1. Notable mention 1 [less than 50 words] \n"
                "2. Notable mention 2 [less than 50 words] \n"
                "# End Example\n\n")
        elif cv_function == "CV Analyzer":
            prompt = """
                    You are evaluating the CV of Software Engineering Candidates.\n                    
                    The ideal candidate will have a strong background in software development.\n                    
                    This individual will work closely with cross-functional teams to design, develop, and maintain innovative software solutions that meet business requirements and align with the company's strategic goals.\n
                    Responsibilities:\n
                    -Design, develop, test, and maintain software applications using modern development methodologies and best practices.\n
                    -Collaborate with product owners, project managers, and business analysts to gather, analyze, and prioritize requirements and convert them into functional specifications.\n
                    -Contribute to the development of application architecture, ensuring scalability, performance, and maintainability.\n 
                    -Integrate software components and third-party applications, ensuring seamless data exchange and system interoperability.\n
                    -Debug and resolve software defects and issues, providing root cause analysis and implementing corrective actions.\n
                    -Participate in code reviews, ensuring adherence to coding standards and best practices.\n
                    -Continuously improve development processes, tools, and methodologies to enhance software quality and team efficiency.\n
                    -Keep up-to-date with emerging technologies, industry trends, and best practices to continuously improve technical skills.\n
                    -Provide technical guidance and support to junior team members as needed.\n
                    -Assist in the creation and maintenance of technical documentation, including design documents, user guides, and system manuals.\n\n
                    Respond with "Qualified" if the candidate meets the above criteria based on their CV and provide your rationale why the candidate is qualified.\n
                    Response with "Not Qualified" if the candidate does not meet the above criteria based on their CV the criteria and provide your rationale why the candidate is not qualified.           

                    """
        elif cv_function == "CV QnA":
            prompt = "Follow the instructions below. Answer the questions based on the contents of the CV.\n"

    return prompt